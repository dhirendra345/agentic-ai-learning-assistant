from pypdf import PdfReader
from pptx import Presentation
from docx import Document


class DocumentLoader:

    @staticmethod
    def read_pdf(file_path):
        text = ""

        pdf = PdfReader(file_path)

        for page in pdf.pages:
            page_text = page.extract_text()

            if page_text:
                text += page_text + "\n"

        return text

    @staticmethod
    def read_ppt(file_path):
        text = ""

        prs = Presentation(file_path)

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

        return text

    @staticmethod
    def read_docx(file_path):
        text = ""

        doc = Document(file_path)

        for para in doc.paragraphs:
            text += para.text + "\n"

        return text